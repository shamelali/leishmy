"""Combine the three Leish! launch plan decks into one PPTX and export all decks to PDF.

Usage: python3 scripts/export-launch-deck-pdfs.py
Outputs (in docs/):
  launch-marketing-plan-all-tiers.pptx   — combined deck (all 3 tiers + dividers)
  launch-marketing-plan.pdf              — High tier
  launch-marketing-plan-medium.pdf       — Medium tier
  launch-marketing-plan-low.pdf          — Low tier
  launch-marketing-plan-all-tiers.pdf    — combined

Requires: python-pptx, reportlab
"""
import os
import re
from copy import deepcopy

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from reportlab.pdfgen import canvas as pdfcanvas
from reportlab.lib.colors import HexColor
from reportlab.pdfbase.pdfmetrics import stringWidth

DOCS = os.path.join(os.path.dirname(__file__), "..", "docs")
EMU_PT = 12700.0


# --------------------------------------------------------------------------
# 1) Combined deck
# --------------------------------------------------------------------------
ROSE = RGBColor(0xE1, 0x1D, 0x48)
PINK = RGBColor(0xDB, 0x27, 0x77)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
DARK2 = RGBColor(0x11, 0x18, 0x27)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x9C, 0xA3, 0xAF)
GREEN = RGBColor(0x65, 0xA3, 0x0D)
SERIF = "Georgia"
SANS = "Calibri"


def _box(s, x, y, w, h, fill, shape=MSO_SHAPE.RECTANGLE, radius=None):
    sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    sp.line.fill.background()
    sp.shadow.inherit = False
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    return sp


def _text(s, x, y, w, h, runs, size=14, color=WHITE, bold=False, font=SANS,
          align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.0):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if isinstance(runs, str):
        runs = [runs]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        if isinstance(para, str):
            para = [(para, {})]
        for t, st in para:
            r = p.add_run()
            r.text = t
            r.font.name = st.get("font", font)
            r.font.size = Pt(st.get("size", size))
            r.font.bold = st.get("bold", bold)
            r.font.italic = st.get("italic", False)
            r.font.color.rgb = st.get("color", color)
    return tb


def _make_blank(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid()
    r.fill.fore_color.rgb = WHITE
    r.line.fill.background()
    r.shadow.inherit = False
    return s


def _copy_slide(prs, src):
    dest = _make_blank(prs)
    for shp in src.shapes:
        dest.shapes._spTree.append(deepcopy(shp._element))
    return dest


def _divider(prs, kicker, title, color, note):
    s = _make_blank(prs)
    _box(s, 0, 0, 13.333, 0.14, color)
    _box(s, 0, 7.36, 13.333, 0.14, color)
    _box(s, 11.2, -1.9, 4.6, 4.6, color, shape=MSO_SHAPE.OVAL)
    _text(s, 0.9, 2.3, 11.0, 0.4, kicker.upper(), size=13, color=color, bold=True)
    _text(s, 0.9, 2.8, 11.5, 1.3, title, size=38, bold=True, font=SERIF)
    _text(s, 0.9, 4.3, 11.0, 0.5, note, size=14, color=MUTED)
    _text(s, 0.9, 6.6, 11.0, 0.4, [[("#BookBeautyAnywhere  ·  #NeverDMForADate", {"size": 12, "color": MUTED, "italic": True})]])
    return s


def merge_decks():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Master title
    s = prs.slides.add_slide(prs.slide_layouts[6])
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid()
    r.fill.fore_color.rgb = DARK2
    r.line.fill.background()
    r.shadow.inherit = False
    _box(s, 0, 0, 13.333, 0.14, ROSE)
    _box(s, 0, 7.36, 13.333, 0.14, PURPLE)
    _box(s, 11.0, -1.7, 4.3, 4.3, ROSE, shape=MSO_SHAPE.OVAL)
    _box(s, -1.5, 5.2, 3.9, 3.9, PURPLE, shape=MSO_SHAPE.OVAL)
    _text(s, 0.9, 1.55, 11.0, 0.4, "LEISH.MY · LAUNCH MARKETING PLAN · 90 DAYS", size=13, color=ROSE, bold=True)
    _text(s, 0.9, 2.05, 11.8, 1.4, [
        [("All Budget Tiers", {"size": 46, "bold": True, "font": SERIF})],
        [("One launch · October 5, 2026 · Three ways to play it", {"size": 18, "color": MUTED})],
    ], line_spacing=1.15)
    tiers = [
        ("HIGH", "RM 96,000", "1,500 clients · 25–30 bookings/wk", ROSE),
        ("MEDIUM", "RM 40,000", "800 clients · 15–20 bookings/wk", PINK),
        ("LOW", "RM 10,000", "300 clients · 6–10 bookings/wk", GREEN),
    ]
    cx = 0.9
    for name, budget, target, c in tiers:
        _box(s, cx, 4.1, 3.7, 1.5, None, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08)
        sp = s.shapes[-1]
        sp.line.color.rgb = c
        sp.line.width = Pt(1.5)
        _text(s, cx + 0.3, 4.3, 3.2, 0.35, name, size=13, color=c, bold=True)
        _text(s, cx + 0.3, 4.65, 3.2, 0.5, budget, size=22, bold=True, font=SERIF)
        _text(s, cx + 0.3, 5.15, 3.2, 0.35, target, size=10.5, color=MUTED)
        cx += 3.9
    _text(s, 0.9, 6.35, 11.5, 0.4, "Full details: docs/launch-marketing-plan*.md · decks: launch-marketing-plan*.pptx", size=11, color=MUTED)

    srcs = [
        ("launch-marketing-plan.pptx", "TIER 1 · HIGH BUDGET", "RM 96,000 — full-scale launch", ROSE),
        ("launch-marketing-plan-medium.pptx", "TIER 2 · MEDIUM BUDGET", "RM 40,000 — lean launch", PINK),
        ("launch-marketing-plan-low.pptx", "TIER 3 · LOW BUDGET", "RM 10,000 — bootstrap proof", GREEN),
    ]
    for fname, kick, title, c in srcs:
        _divider(prs, kick, title, c, "Same launch date · same campaign · scaled execution — see docs/launch-marketing-plan-tiers.md for the comparison.")
        src = Presentation(os.path.join(DOCS, fname))
        for slide in src.slides:
            _copy_slide(prs, slide)

    # Closing
    s = prs.slides.add_slide(prs.slide_layouts[6])
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid()
    r.fill.fore_color.rgb = DARK2
    r.line.fill.background()
    r.shadow.inherit = False
    _box(s, 0, 0, 13.333, 0.14, ROSE)
    _box(s, 0, 7.36, 13.333, 0.14, PURPLE)
    _box(s, 11.0, -1.7, 4.3, 4.3, PURPLE, shape=MSO_SHAPE.OVAL)
    _text(s, 0.9, 2.4, 11.5, 1.2, "Pick a tier. Launch Oct 5.", size=44, bold=True, font=SERIF)
    _text(s, 0.9, 3.7, 11.0, 0.5, [
        [("Decision path: ", {"size": 14, "bold": True, "color": ROSE}),
         ("Low proves economics → Medium grows it → High captures wedding season. Close the 7 blockers by Aug 10.", {"size": 14, "color": MUTED})],
    ])
    _text(s, 0.9, 6.3, 11.0, 0.4, [[("#BookBeautyAnywhere  ·  #NeverDMForADate  ·  Your Beauty, Perfected.", {"size": 13, "color": MUTED, "italic": True})]])

    out = os.path.join(DOCS, "launch-marketing-plan-all-tiers.pptx")
    prs.save(out)
    print(f"Saved: {out} ({len(prs.slides._sldIdLst)} slides)")
    return out


# --------------------------------------------------------------------------
# 2) PDF renderer (approximate faithful render of the generated decks)
# --------------------------------------------------------------------------
FONT_MAP = {
    "Georgia": ("Times-Roman", "Times-Bold", "Times-Italic", "Times-BoldItalic"),
    "Calibri": ("Helvetica", "Helvetica-Bold", "Helvetica-Oblique", "Helvetica-BoldOblique"),
    "Consolas": ("Courier", "Courier-Bold", "Courier-Oblique", "Courier-BoldOblique"),
}


def _hex(c):
    return HexColor("#%02X%02X%02X" % (c[0], c[1], c[2]))


def _runs_of(para, default_size, default_color):
    out = []
    for run in para.runs:
        t = run.text
        if not t:
            continue
        name = run.font.name or "Calibri"
        base, b, i, bi = FONT_MAP.get(name, FONT_MAP["Calibri"])
        bold = bool(run.font.bold)
        italic = bool(run.font.italic)
        font = bi if (bold and italic) else (b if bold else (i if italic else base))
        size = run.font.size.pt if run.font.size else default_size
        col = run.font.color.rgb if run.font.color and run.font.color.type is not None else default_color
        out.append((t, font, size, _hex(col)))
    return out


def _wrap_tokens(tokens, width):
    """Greedy word-wrap over mixed-format tokens. Returns list of (tokens, line_width)."""
    lines, cur, cur_w = [], [], 0.0
    space_w = stringWidth(" ", "Helvetica", 10)  # placeholder; recomputed per font below
    for text, font, size, col in tokens:
        words = re.split(r"(\s+)", text)
        for w in words:
            if w == "":
                continue
            ww = stringWidth(w, font, size)
            if w.isspace():
                if cur:
                    cur.append((w, font, size, col))
                    cur_w += ww
                continue
            if cur and cur_w + ww > width:
                lines.append(cur)
                cur, cur_w = [], 0.0
            cur.append((w, font, size, col))
            cur_w += ww
    if cur:
        lines.append(cur)
    return lines


def _para_lines(para, box_w, default_size, default_color):
    tokens = _runs_of(para, default_size, default_color)
    if not tokens:
        return [], 0
    lines = _wrap_tokens(tokens, box_w)
    return lines, max((sz for _, _, sz, _ in tokens))


def _render_text(c, shp):
    tf = shp.text_frame
    x = shp.left / EMU_PT
    y_top = shp.top / EMU_PT
    w, h = shp.width / EMU_PT, shp.height / EMU_PT
    ml = tf.margin_left / EMU_PT if tf.margin_left else 0
    mr = tf.margin_right / EMU_PT if tf.margin_right else 0
    mt = tf.margin_top / EMU_PT if tf.margin_top else 0
    mb = tf.margin_bottom / EMU_PT if tf.margin_bottom else 0
    avail_w = max(w - ml - mr, 10)
    avail_h = max(h - mt - mb, 0)
    default_color = RGBColor(0x1F, 0x29, 0x37)
    default_size = 18.0

    blocks = []  # (align, lines, line_h, font_heights)
    for para in tf.paragraphs:
        lines, max_sz = _para_lines(para, avail_w, default_size, default_color)
        if not lines:
            blocks.append((para.alignment, [], 0, 0))
            continue
        ls = para.line_spacing
        if isinstance(ls, (int, float)):
            line_h = max_sz * 1.2 * ls
        else:
            line_h = (ls / EMU_PT) if ls else max_sz * 1.2
        blocks.append((para.alignment, lines, line_h, max_sz))

    total_h = sum(b[2] * len(b[1]) for b in blocks if b[1])
    anchor = tf.vertical_anchor
    flipped_top = 540 - y_top - h
    if anchor == MSO_ANCHOR.MIDDLE:
        y0 = flipped_top + mt + max((avail_h - total_h) / 2, 0)
    elif anchor == MSO_ANCHOR.BOTTOM:
        y0 = flipped_top + mt + max(avail_h - total_h, 0)
    else:
        y0 = flipped_top + mt

    ty = y0
    for align, lines, line_h, max_sz in blocks:
        for toks in lines:
            text_w = sum(stringWidth(t, f, s) for t, f, s, _ in toks)
            if align == PP_ALIGN.CENTER:
                tx = x + ml + (avail_w - text_w) / 2
            elif align == PP_ALIGN.RIGHT:
                tx = x + ml + (avail_w - text_w)
            else:
                tx = x + ml
            for t, f, s, col in toks:
                c.setFont(f, s)
                c.setFillColor(col)
                c.drawString(tx, ty + (max_sz - s) * 0.3, t)
                tx += stringWidth(t, f, s)
            ty += line_h


def _render_autoshape(c, shp):
    x, y = shp.left / EMU_PT, (540 - (shp.top / EMU_PT) - (shp.height / EMU_PT))
    w, h = shp.width / EMU_PT, shp.height / EMU_PT
    if w <= 0 or h <= 0:
        return
    fill = None
    if shp.fill.type == 1:  # solid
        fill = _hex((shp.fill.fore_color.rgb[0], shp.fill.fore_color.rgb[1], shp.fill.fore_color.rgb[2]))
    stroke = None
    sw = 0
    if shp.line.fill.type == 1:
        stroke = _hex((shp.line.color.rgb[0], shp.line.color.rgb[1], shp.line.color.rgb[2]))
        sw = shp.line.width.pt if shp.line.width else 1.0
    st = shp.auto_shape_type
    if fill is None:
        c.setFillColor(HexColor("#000000"))
    else:
        c.setFillColor(fill)
    if stroke is None:
        c.setStrokeColor(HexColor("#000000"))
    else:
        c.setStrokeColor(stroke)
    c.setLineWidth(sw)
    if st == MSO_SHAPE.OVAL:
        c.ellipse(x, y, x + w, y + h, stroke=stroke is not None, fill=fill is not None)
    elif st == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            adj = shp.adjustments[0]
        except Exception:
            adj = 0.16667
        radius = min(adj * min(w, h), min(w, h) / 2)
        c.roundRect(x, y, w, h, radius, stroke=stroke is not None, fill=fill is not None)
    else:
        c.rect(x, y, w, h, stroke=stroke is not None, fill=fill is not None)


def _render_table(c, shp):
    tbl = shp.table
    x0, y_top = shp.left / EMU_PT, shp.top / EMU_PT
    cols = []
    for col in tbl.columns:
        cols.append(col.width / EMU_PT)
    rows = []
    for row in tbl.rows:
        rows.append(row.height / EMU_PT)
    y = 540 - y_top
    for ri, row in enumerate(tbl.rows):
        x = x0
        rh = rows[ri]
        y -= rh
        for ci, cell in enumerate(row.cells):
            cw = cols[ci]
            fill = None
            try:
                if cell.fill.type == 1:
                    rgb = cell.fill.fore_color.rgb
                    fill = _hex((rgb[0], rgb[1], rgb[2]))
            except Exception:
                pass
            if fill:
                c.setFillColor(fill)
                c.rect(x, y, cw, rh, stroke=0, fill=1)
            ml = (cell.margin_left or 0) / EMU_PT
            mr = (cell.margin_right or 0) / EMU_PT
            mt = (cell.margin_top or 0) / EMU_PT
            mb = (cell.margin_bottom or 0) / EMU_PT
            avail_w = max(cw - ml - mr, 8)
            avail_h = max(rh - mt - mb, 0)
            first = cell.text_frame.paragraphs[0]
            default_color = RGBColor(0x1F, 0x29, 0x37)
            blocks = []
            for para in cell.text_frame.paragraphs:
                lines, max_sz = _para_lines(para, avail_w, 12, default_color)
                ls = para.line_spacing
                line_h = max_sz * 1.2 * ls if isinstance(ls, (int, float)) else max_sz * 1.2
                blocks.append((para.alignment, lines, line_h, max_sz))
            total_h = sum(b[2] * len(b[1]) for b in blocks if b[1])
            anchor = cell.vertical_anchor
            if anchor == MSO_ANCHOR.MIDDLE:
                ty = y + (rh - total_h) / 2
            elif anchor == MSO_ANCHOR.BOTTOM:
                ty = y + rh - total_h - mb
            else:
                ty = y + mt
            for align, lines, line_h, max_sz in blocks:
                for toks in lines:
                    text_w = sum(stringWidth(t, f, s) for t, f, s, _ in toks)
                    if align == PP_ALIGN.CENTER:
                        tx = x + ml + (avail_w - text_w) / 2
                    elif align == PP_ALIGN.RIGHT:
                        tx = x + ml + avail_w - text_w
                    else:
                        tx = x + ml
                    for t, f, s, col in toks:
                        c.setFont(f, s)
                        c.setFillColor(col)
                        c.drawString(tx, ty + (max_sz - s) * 0.3, t)
                        tx += stringWidth(t, f, s)
                    ty += line_h
            x += cw


def render_pdf(pptx_path, pdf_path):
    prs = Presentation(pptx_path)
    c = pdfcanvas.Canvas(pdf_path, pagesize=(960, 540))
    for slide in prs.slides:
        for shp in slide.shapes:
            try:
                if shp.has_table:  # TABLE graphicFrame
                    _render_table(c, shp)
                elif shp.shape_type == 1:  # AUTO_SHAPE (fill first — chips also carry text)
                    _render_autoshape(c, shp)
                elif shp.has_text_frame:  # TEXT_BOX
                    _render_text(c, shp)
            except Exception as e:
                print(f"  [warn] shape skipped on slide: {e}")
        c.showPage()
    c.save()
    print(f"Saved: {pdf_path} ({len(prs.slides._sldIdLst)} pages)")


def pdf_page_count(path):
    with open(path, "rb") as f:
        data = f.read()
    return len(re.findall(rb"/Type\s*/Page[^s]", data))


if __name__ == "__main__":
    combined = merge_decks()
    decks = [
        os.path.join(DOCS, "launch-marketing-plan.pptx"),
        os.path.join(DOCS, "launch-marketing-plan-medium.pptx"),
        os.path.join(DOCS, "launch-marketing-plan-low.pptx"),
        combined,
    ]
    for d in decks:
        pdf = os.path.splitext(d)[0] + ".pdf"
        print("Rendering:", os.path.basename(d))
        render_pdf(d, pdf)
        print(f"  pages verified: {pdf_page_count(pdf)}")
    print("Done.")
