"""Generate the Leish! launch marketing plan slide deck (docs/launch-marketing-plan.pptx).

Usage: python3 scripts/generate-launch-deck.py
Regenerate anytime the plan in docs/launch-marketing-plan.md changes.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# ---- Brand palette (matches leish.my site) ----
ROSE = RGBColor(0xE1, 0x1D, 0x48)     # rose-600
PINK = RGBColor(0xDB, 0x27, 0x77)     # pink-600
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)   # purple-500
DARK = RGBColor(0x1F, 0x29, 0x37)     # gray-800
DARK2 = RGBColor(0x11, 0x18, 0x27)
GRAY = RGBColor(0x6B, 0x72, 0x80)     # gray-500
LIGHT = RGBColor(0xFD, 0xF2, 0xF8)    # rose-50
LIGHT2 = RGBColor(0xF5, 0xF3, 0xFF)   # violet-50
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x9C, 0xA3, 0xAF)

SERIF = "Georgia"
SANS = "Calibri"

OUT = os.path.join(os.path.dirname(__file__), "..", "docs", "launch-marketing-plan.pptx")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height
PAGE = [0]


def slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid()
    r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    PAGE[0] += 1
    return s


def box(s, x, y, w, h, fill=None, line=None, shape=MSO_SHAPE.RECTANGLE, radius=None, shadow=False):
    sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    sp.shadow.inherit = shadow
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    return sp


def text(s, x, y, w, h, runs, size=14, color=DARK, bold=False, font=SANS,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.0, space_after=0):
    """runs: str OR list of paragraphs; each paragraph is str or list of (text, dict) runs."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if isinstance(runs, str):
        runs = [runs]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        if space_after:
            p.space_after = Pt(space_after)
        if isinstance(para, str):
            para = [(para, {})]
        for t, st in para:
            r = p.add_run()
            r.text = t
            r.font.name = st.get("font", font)
            r.font.size = Pt(st.get("size", size))
            r.font.bold = st.get("bold", bold)
            r.font.italic = st.get("italic", False)
            r.font.color.rgb = st.get("color", color)
    return tb


def footer(s, note="Leish! Launch Marketing Plan · Oct 5, 2026 · RM 96,000 / 90 days"):
    box(s, 0, 7.12, SW, 0.38, fill=WHITE)
    text(s, 0.6, 7.14, 9.0, 0.3, note, size=9, color=MUTED)
    text(s, 12.3, 7.14, 0.5, 0.3, str(PAGE[0]), size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def header(s, kicker, title, accent=ROSE):
    text(s, 0.6, 0.42, 11.5, 0.3, kicker.upper(), size=11, color=accent, bold=True)
    text(s, 0.6, 0.72, 12.1, 0.8, title, size=30, color=DARK, bold=True, font=SERIF)
    box(s, 0.62, 1.52, 1.6, 0.045, fill=accent)


def chip_row(s, y, items, x=0.6, color=ROSE, w=2.1, h=0.4, gap=0.15, size=12, bg=LIGHT):
    cx = x
    for it in items:
        box(s, cx, y, w, h, fill=bg, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
        text(s, cx, y + 0.045, w, h - 0.09, it, size=size, color=color, bold=True, align=PP_ALIGN.CENTER)
        cx += w + gap


# ============================================================ 1. TITLE
s = slide(DARK2)
box(s, 0, 0, SW, 0.14, fill=ROSE)
box(s, 0, 7.36, SW, 0.14, fill=PURPLE)
box(s, 10.9, -1.6, 4.4, 4.4, fill=ROSE, shape=MSO_SHAPE.OVAL)
box(s, 11.7, -0.9, 2.6, 2.6, fill=PINK, shape=MSO_SHAPE.OVAL)
box(s, -1.4, 5.4, 3.8, 3.8, fill=PURPLE, shape=MSO_SHAPE.OVAL)
box(s, -0.7, 6.1, 2.2, 2.2, fill=ROSE, shape=MSO_SHAPE.OVAL)
text(s, 0.9, 1.55, 9.5, 0.4, "LEISH.MY · LAUNCH MARKETING PLAN · 90 DAYS", size=13, color=ROSE, bold=True)
text(s, 0.9, 2.05, 11.5, 1.9, [
    [("Never DM for a Date Again.", {"size": 44, "bold": True, "font": SERIF, "color": WHITE})],
    [("Malaysia's beauty booking marketplace — official launch ", {"size": 20, "color": MUTED}),
     ("October 5, 2026", {"size": 20, "bold": True, "color": WHITE})],
], line_spacing=1.1)
box(s, 0.95, 4.15, 4.4, 1.15, fill=None, line=ROSE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
text(s, 1.2, 4.32, 4.0, 0.9, [
    [("TOTAL BUDGET", {"size": 11, "color": MUTED, "bold": True})],
    [("RM 96,000", {"size": 26, "bold": True, "color": WHITE})],
], line_spacing=1.05)
box(s, 5.6, 4.15, 4.4, 1.15, fill=None, line=PINK, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
text(s, 5.85, 4.32, 4.0, 0.9, [
    [("DAY-90 TARGETS", {"size": 11, "color": MUTED, "bold": True})],
    [("25–30 bookings/wk · 60+ artists · 1,500 clients", {"size": 15, "bold": True, "color": WHITE})],
], line_spacing=1.05)
box(s, 10.25, 4.15, 2.45, 1.15, fill=None, line=PURPLE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
text(s, 10.5, 4.32, 2.0, 0.9, [
    [("LAUNCH DATE", {"size": 11, "color": MUTED, "bold": True})],
    [("Oct 5, 2026", {"size": 20, "bold": True, "color": WHITE})],
], line_spacing=1.05)
text(s, 0.9, 5.75, 11.0, 0.9, [
    [("Aligned with the FY26 financial forecast — Year-1 revenue RM 180–250k · 20% commission · avg booking RM 250–350", {"size": 13, "color": MUTED})],
], line_spacing=1.2)

# ============================================================ 2. THE PLAY
s = slide()
header(s, "Executive summary", "The Play")
cols = [
    ("WHY", ROSE, "Beauty booking in Malaysia is broken",
     "Clients DM 15 artists and wait days for a quote. Artists drown in DM chaos, ghosting and late payments. Nobody owns real-time MUA booking — that's Leish!'s wedge."),
    ("WHAT", PINK, "A two-sided marketplace",
     "Clients book vetted makeup artists & studios with real-time availability, instant confirmation and secure Billplz payment. Artists get portfolios, bookings, reminders and payouts."),
    ("HOW", PURPLE, "Both sides, supply first",
     "Recruit 30 anchor artists before Day 0 (critical path), then drive client demand through paid, KOLs, events and PR. Launch at the peak of wedding season."),
]
cx = 0.6
for title, c, head, body in cols:
    box(s, cx, 2.0, 3.9, 4.3, fill=LIGHT if c != PURPLE else LIGHT2, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
    box(s, cx, 2.0, 3.9, 0.14, fill=c)
    text(s, cx + 0.35, 2.45, 3.2, 0.3, title, size=13, color=c, bold=True)
    text(s, cx + 0.35, 2.85, 3.2, 0.9, head, size=19, bold=True, font=SERIF, color=DARK)
    text(s, cx + 0.35, 3.75, 3.25, 2.4, body, size=13, color=GRAY, line_spacing=1.2)
    cx += 4.05
chip_row(s, 6.65, ["Launch theme: Never DM for a date again", "Artist theme: Get booked. Get paid.", "Hero categories: Bridal · Hijab · Event"], w=3.92, h=0.42, size=11.5)

# ============================================================ 3. POSITIONING
s = slide()
header(s, "Positioning", "Two messages, one marketplace")
box(s, 0.6, 2.0, 6.0, 2.7, fill=LIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.07)
box(s, 0.6, 2.0, 0.12, 2.7, fill=ROSE)
text(s, 0.95, 2.3, 5.3, 0.35, "TO CLIENTS", size=12, color=ROSE, bold=True)
text(s, 0.95, 2.7, 5.4, 1.8, [
    [("“Book Malaysia's top makeup artists & studios in minutes — real-time availability, instant confirmation, secure payment.”", {"size": 17, "bold": True, "font": SERIF})],
], line_spacing=1.15)
box(s, 6.85, 2.0, 6.0, 2.7, fill=LIGHT2, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.07)
box(s, 6.85, 2.0, 0.12, 2.7, fill=PURPLE)
text(s, 7.2, 2.3, 5.3, 0.35, "TO ARTISTS & STUDIOS", size=12, color=PURPLE, bold=True)
text(s, 7.2, 2.7, 5.4, 1.8, [
    [("“Get booked. Get paid. Your portfolio, your pricing, zero admin — we handle bookings, reminders and payments.”", {"size": 17, "bold": True, "font": SERIF})],
], line_spacing=1.15)
text(s, 0.6, 5.0, 12.0, 0.35, "WHAT THE PLATFORM ALREADY HAS — USE IT", size=12, color=GRAY, bold=True)
chips = ["Instant booking + Billplz payments", "Real-time availability & reviews", "Artist/studio profiles with portfolios", "Rewards + Leish+ (RM29/mo) upsell", "WhatsApp booking flow", "Brevo email (hello@leish.my)", "GA4 + Meta Pixel installed", "Artist shareable profile links"]
cx, cy = 0.6, 5.4
for i, c in enumerate(chips):
    box(s, cx, cy, 3.05, 0.44, fill=WHITE, line=RGBColor(0xFD, 0xD1, 0xDD) if i % 2 == 0 else RGBColor(0xDD, 0xD6, 0xFE), shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    text(s, cx + 0.12, cy + 0.05, 2.85, 0.36, c, size=11, color=DARK, align=PP_ALIGN.CENTER, bold=True)
    cx += 3.17
    if cx > 10.0:
        cx = 0.6
        cy += 0.58

# ============================================================ 4. TIMELINE
s = slide()
header(s, "Timeline", "Four phases, one launch")
phases = [
    ("PHASE 0 · PRE-LAUNCH", "Aug 10 – Oct 4", ROSE,
     ["Waitlist live with RM30 offer", "Recruit 30 anchor artists (concierge)", "Teaser social 3x/wk + countdown", "GA4/Pixel events + Brevo sequences"]),
    ("PHASE 1 · LAUNCH WEEK", "Oct 5 – 11", PINK,
     ["“Never DM for a date again” campaign", "KL launch pop-up + press event", "Wave 1: 8–12 micro KOLs", "Meta + Google paid ON, PR out"]),
    ("PHASE 2 · WEDDING SEASON", "Oct 12 – Dec 31", PURPLE,
     ["Wave 2 KOLs + client UGC", "Referral program live (Day 30)", "Wedding expos + hijab push", "B2B corporate + Leish+ pilot"]),
    ("PHASE 3 · REVIEW", "Jan 2027", RGBColor(0x0E, 0x94, 0x92), ["Score vs Day-90 targets", "Cohort/LTV analysis", "Q1'27 plan + expansion cities"]),
]
cx = 0.6
for title, dates, c, items in phases:
    box(s, cx, 2.05, 3.0, 4.35, fill=WHITE, line=RGBColor(0xE5, 0xE7, 0xEB), shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
    box(s, cx, 2.05, 3.0, 0.14, fill=c)
    text(s, cx + 0.25, 2.4, 2.55, 0.55, title, size=13, color=c, bold=True)
    text(s, cx + 0.25, 2.95, 2.55, 0.3, dates, size=12, color=GRAY, bold=True)
    yy = 3.4
    for it in items:
        box(s, cx + 0.28, yy + 0.07, 0.09, 0.09, fill=c, shape=MSO_SHAPE.OVAL)
        text(s, cx + 0.48, yy, 2.35, 0.9, it, size=11.5, color=DARK, line_spacing=1.05)
        yy += 0.78
    cx += 3.13

# ============================================================ 5. KPIs
s = slide()
header(s, "North-star metric: completed bookings per week", "Day-90 Targets")
stats = [
    ("1,500", "Registered clients", ROSE),
    ("60+", "Active artists & studios", PINK),
    ("25–30", "Bookings / week", PURPLE),
    ("RM 7–9k", "GMV / week", RGBColor(0x0E, 0x94, 0x92)),
    ("≤ RM 20", "Client CAC (paid)", RGBColor(0xD9, 0x77, 0x06)),
    ("2.5–3.5x", "Blended ROAS", RGBColor(0x7C, 0x3A, 0xED)),
]
cx, cy = 0.6, 2.0
for val, label, c in stats:
    box(s, cx, cy, 3.95, 1.35, fill=LIGHT if c != PURPLE else LIGHT2, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08)
    box(s, cx, cy, 0.1, 1.35, fill=c)
    text(s, cx + 0.35, cy + 0.18, 3.4, 0.6, val, size=27, bold=True, font=SERIF, color=c)
    text(s, cx + 0.35, cy + 0.85, 3.4, 0.4, label, size=12.5, color=GRAY, bold=True)
    cx += 4.08
    if cx > 12.0:
        cx = 0.6
        cy += 1.5
rows = [
    ("KPI", "Pre-launch", "Day 30", "Day 60", "Day 90"),
    ("Waitlist / clients", "2,000 waitlist", "600", "1,000", "1,500"),
    ("Weekly bookings", "—", "12", "20", "25–30"),
    ("Commission revenue / wk (20%)", "—", "RM 700", "RM 1,200", "RM 1,500–1,800"),
    ("Email list", "2,500", "4,000", "5,000", "6,000"),
    ("Social following", "3,000", "6,000", "8,000", "10,000"),
]
tbl = s.shapes.add_table(len(rows), 5, Inches(0.6), Inches(5.35), Inches(12.13), Inches(1.55)).table
tbl.columns[0].width = Inches(4.4)
for i in range(1, 5):
    tbl.columns[i].width = Inches(1.93)
for ri, row in enumerate(rows):
    for ci, val in enumerate(row):
        cell = tbl.cell(ri, ci)
        cell.margin_left = cell.margin_right = Inches(0.12)
        cell.margin_top = cell.margin_bottom = Inches(0.03)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK if ri == 0 else (WHITE if ri % 2 else RGBColor(0xFB, 0xF7, 0xFB))
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER if ci else PP_ALIGN.LEFT
        r = p.add_run()
        r.text = val
        r.font.size = Pt(12 if ri == 0 else 11.5)
        r.font.bold = ri == 0 or ci == 0
        r.font.color.rgb = WHITE if ri == 0 else DARK
        r.font.name = SANS

# ============================================================ 6. BUDGET
s = slide()
header(s, "RM 96,000 · 90 days", "Budget Breakdown")
items = [
    ("Meta ads (IG + FB, incl. retargeting)", 29000, ROSE, 0.30),
    ("Influencers / KOLs (2 waves, 16–18 creators)", 24000, PINK, 0.25),
    ("Google Search + PMax", 12000, PURPLE, 0.13),
    ("Events — launch pop-up + wedding expos", 12000, RGBColor(0xD9, 0x77, 0x06), 0.13),
    ("TikTok ads (Spark on creator content)", 6000, RGBColor(0x0E, 0x94, 0x92), 0.06),
    ("Artist acquisition (ads + webinar)", 6000, RGBColor(0x7C, 0x3A, 0xED), 0.06),
    ("Community / campus seeding", 4000, RGBColor(0x65, 0xA3, 0x0D), 0.04),
    ("PR & media kit", 3000, RGBColor(0x94, 0x78, 0x55), 0.03),
]
yy = 2.0
for name, amt, c, pct in items:
    text(s, 0.6, yy + 0.03, 5.6, 0.35, f"{name}", size=12.5, color=DARK, bold=True)
    text(s, 6.35, yy + 0.03, 1.6, 0.35, f"RM {amt:,}", size=12.5, color=GRAY, bold=True, align=PP_ALIGN.RIGHT)
    box(s, 8.1, yy + 0.09, 3.1, 0.24, fill=RGBColor(0xF1, 0xF2, 0xF6), shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    box(s, 8.1, yy + 0.09, 3.1 * pct, 0.24, fill=c, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    text(s, 11.35, yy + 0.03, 1.4, 0.35, f"{int(pct*100)}%", size=12.5, color=DARK, bold=True, align=PP_ALIGN.RIGHT)
    yy += 0.56
box(s, 0.6, 6.55, 12.1, 0.02, fill=RGBColor(0xE5, 0xE7, 0xEB))
text(s, 0.6, 6.7, 12.1, 0.4, [
    [("Payback logic: ", {"size": 12, "bold": True, "color": DARK}),
     ("blended CAC RM20 → client LTV RM 90–150 (2–3 bookings × RM 300 × 20%). CAC pays back in 1–2 bookings. Kill or scale by ROAS weekly — this is a growth test, not a burn.", {"size": 12, "color": GRAY})],
], line_spacing=1.15)

# ============================================================ 7. CHANNELS
s = slide()
header(s, "Channel mix", "Owned · Paid · Earned")
groups = [
    ("OWNED — foundation", ROSE, [
        "Instagram — 4–5 posts + 10 stories + 2 reels / wk",
        "TikTok — 3–4 UGC-style reels / wk",
        "Facebook + bridal/beauty groups — daily participation",
        "Email (Brevo) — lifecycle sequences + 2 newsletters/mo",
        "WhatsApp API — confirmations, reminders, broadcasts",
        "SEO — 2 blogs/mo on /inspiration + /events",
        "Artist share links — every artist recruits clients",
    ]),
    ("PAID — RM 47k", PINK, [
        "Meta (IG+FB) RM 24k — bridal/beauty interest stacks, KL-Selangor, 20–40",
        "Google RM 12k — high-intent: “bridal makeup KL”, “MUA near me”",
        "TikTok RM 6k — Spark Ads on creator content",
        "Retargeting RM 5k — visitors, waitlist & cart abandoners, RM30 creative",
        "Artist-side ads — target working MUAs (“switch from DM chaos”)",
    ]),
    ("EARNED + COMMUNITY — RM 43k", PURPLE, [
        "2 KOL waves — 16–18 micro/mid creators, flat + booking-code commission",
        "1–2 wedding expos + KL launch pop-up (glam bar)",
        "PR — SAYS, The Star, Malay Mail, NST, wedding blogs",
        "Giveaways in bridal FB groups",
        "Campus pop-ups — Cyberjaya (MMU, Limkokwing, UPM area)",
        "Free “grow your MUA business” webinar (Nov)",
    ]),
]
cx = 0.6
for title, c, items in groups:
    box(s, cx, 2.0, 3.98, 4.75, fill=WHITE, line=RGBColor(0xE5, 0xE7, 0xEB), shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
    box(s, cx, 2.0, 3.98, 0.14, fill=c)
    text(s, cx + 0.3, 2.35, 3.4, 0.35, title, size=15, color=c, bold=True, font=SERIF)
    yy = 2.85
    for it in items:
        box(s, cx + 0.32, yy + 0.08, 0.09, 0.09, fill=c, shape=MSO_SHAPE.OVAL)
        text(s, cx + 0.52, yy, 3.25, 0.75, it, size=11, color=DARK, line_spacing=1.05)
        yy += 0.62
    cx += 4.13

# ============================================================ 8. OFFERS
s = slide()
header(s, "Lock before Day 0", "Launch Offers & Promotions")
box(s, 0.6, 2.0, 6.0, 4.4, fill=LIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
box(s, 0.6, 2.0, 6.0, 0.14, fill=ROSE)
text(s, 0.95, 2.3, 5.3, 0.35, "CLIENT SIDE (DEMAND)", size=14, color=ROSE, bold=True)
client_offers = [
    ("LEISH30", "RM30 off first booking · cap 300 · valid to Day 30"),
    ("Early-bird bonus", "Waitlist members get RM30 + priority access"),
    ("Refer-a-friend", "RM20 credit each · live from Day 30"),
]
yy = 2.85
for name, desc in client_offers:
    box(s, 0.95, yy, 5.3, 0.95, fill=WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.1)
    text(s, 1.2, yy + 0.12, 4.9, 0.35, name, size=13, bold=True, color=ROSE)
    text(s, 1.2, yy + 0.5, 4.9, 0.4, desc, size=11.5, color=GRAY)
    yy += 1.1
box(s, 6.85, 2.0, 6.0, 4.4, fill=LIGHT2, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
box(s, 6.85, 2.0, 6.0, 0.14, fill=PURPLE)
text(s, 7.2, 2.3, 5.3, 0.35, "ARTIST SIDE (SUPPLY)", size=14, color=PURPLE, bold=True)
artist_offers = [
    ("Free 3-month Pro listing", "Normally RM 200–500/mo · cap 50 artists"),
    ("0% commission, first 10 bookings", "vs standard 20% · cap 30 artists"),
    ("Artist referral", "RM100 + featured slot when referee hits 5 bookings"),
]
yy = 2.85
for name, desc in artist_offers:
    box(s, 7.2, yy, 5.3, 0.95, fill=WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.1)
    text(s, 7.45, yy + 0.12, 4.9, 0.35, name, size=13, bold=True, color=PURPLE)
    text(s, 7.45, yy + 0.5, 4.9, 0.4, desc, size=11.5, color=GRAY)
    yy += 1.1
text(s, 0.6, 6.55, 12.1, 0.45, [
    [("⚠ Pre-launch dev task: ", {"size": 12, "bold": True, "color": DARK}),
     ("verify promo-code / credit support in admin. If unsupported, use UTM + support-issued credits and log a dev task before Aug 24.", {"size": 12, "color": GRAY})],
], line_spacing=1.15)

# ============================================================ 9. CONTENT CALENDAR
s = slide()
header(s, "First 12 weeks", "Content Calendar — Highlights")
rows = [
    ("Week", "Theme", "Hero content", "Campaign"),
    ("W1 · Aug 10", "Tease", "Founder story reel — “Why we built Leish!”", "Waitlist open"),
    ("W2", "Artist spotlight", "“Meet the artist” carousels ×3", "Artist recruitment"),
    ("W3", "Pain-point comedy", "“POV: DM-ing 15 MUAs for your wedding date”", "Waitlist"),
    ("W4", "Education", "Bridal makeup cost guide KL 2026", "Waitlist"),
    ("W5", "Proof", "Beta client before/after + testimonial", "Waitlist"),
    ("W6", "Countdown", "“Leish! launches Oct 5” series", "Launch hype"),
    ("W7 · Oct 5", "LAUNCH WEEK", "Event content, KOL posts, LEISH30 code", "DAY 0"),
    ("W8", "Social proof", "UGC roundup + first bookings", "Post-launch"),
    ("W9", "Hijab series", "Hijab MUA directory + looks", "Community"),
    ("W10", "Wedding season", "“Book your wedding MUA now” urgency", "Paid push"),
    ("W11", "Referrals", "Refer-a-friend launch creative", "Referral"),
    ("W12", "Festive prep", "CNY / Raya early-bird booking", "Festive"),
]
tbl = s.shapes.add_table(len(rows), 4, Inches(0.6), Inches(1.95), Inches(12.13), Inches(4.6)).table
tbl.columns[0].width = Inches(1.7)
tbl.columns[1].width = Inches(2.1)
tbl.columns[2].width = Inches(5.33)
tbl.columns[3].width = Inches(3.0)
for ri, row in enumerate(rows):
    for ci, val in enumerate(row):
        cell = tbl.cell(ri, ci)
        cell.margin_left = cell.margin_right = Inches(0.1)
        cell.margin_top = cell.margin_bottom = Inches(0.02)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        if ri == 0:
            cell.fill.fore_color.rgb = DARK
        elif ri == 7:
            cell.fill.fore_color.rgb = ROSE
        else:
            cell.fill.fore_color.rgb = WHITE if ri % 2 else RGBColor(0xFB, 0xF7, 0xFB)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = val
        r.font.size = Pt(11.5 if ri else 12)
        r.font.bold = ri == 0 or ri == 7 or ci == 0
        r.font.color.rgb = WHITE if ri == 0 or ri == 7 else DARK
        r.font.name = SANS
text(s, 0.6, 6.7, 12.1, 0.4, [
    [("Evergreen SEO (2/mo): ", {"size": 12, "bold": True, "color": DARK}),
     ("“Best Bridal Makeup Artists in KL & Selangor 2026” · “How Much Does Bridal Makeup Cost in Malaysia?” · “Hijab Makeup Looks for Weddings” · “10 Questions to Ask Your MUA”", {"size": 12, "color": GRAY})],
], line_spacing=1.15)

# ============================================================ 10. RISKS
s = slide()
header(s, "Watch weekly", "Risks & Mitigations")
rows = [
    ("Risk", "L", "I", "Mitigation"),
    ("Not enough artists at launch", "Med", "High", "Concierge recruitment from T-6; 30-artist gate — delay Day 0 if short"),
    ("Demand lags supply (artists churn)", "Med", "High", "Weekly zero-booking check; demand ads + featured rotation; commission-free period"),
    ("Checkout friction (Billplz)", "Med", "High", "Track booking_started→completed; A/B deposit vs full payment; WhatsApp assist"),
    ("KOL underperformance", "Med", "Med", "60/40 flat + booking commission; test 4 in wave 1 before wave-2 spend"),
    ("Promo abuse (RM30 codes)", "Med", "Low", "One code per account; cap 300; fraud flag on repeat emails"),
    ("SEO ramp too slow", "High", "Low", "Content engine from Day 1; treat SEO as 6-month play; rely on paid + social at launch"),
    ("Traffic spike / infra", "Low", "High", "Launch-ready ✅; load-test /artists + /bookings; UPSTASH Redis in prod"),
]
tbl = s.shapes.add_table(len(rows), 4, Inches(0.6), Inches(1.95), Inches(12.13), Inches(4.7)).table
tbl.columns[0].width = Inches(3.3)
tbl.columns[1].width = Inches(0.7)
tbl.columns[2].width = Inches(0.7)
tbl.columns[3].width = Inches(7.43)
for ri, row in enumerate(rows):
    for ci, val in enumerate(row):
        cell = tbl.cell(ri, ci)
        cell.margin_left = cell.margin_right = Inches(0.1)
        cell.margin_top = cell.margin_bottom = Inches(0.02)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK if ri == 0 else (WHITE if ri % 2 else RGBColor(0xFB, 0xF7, 0xFB))
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = val
        r.font.size = Pt(12 if ri == 0 else 11.5)
        r.font.bold = ri == 0 or ci == 0
        r.font.color.rgb = WHITE if ri == 0 else DARK
        r.font.name = SANS
text(s, 0.6, 6.8, 12.1, 0.35, [
    [("Guardrails: ", {"size": 12, "bold": True, "color": DARK}),
     ("artists idle 14+ days · no-show/cancellation <15% · artist response <4 hrs · search→booking ≥2%", {"size": 12, "color": GRAY})],
], line_spacing=1.15)

# ============================================================ 11. CHECKLIST
s = slide()
header(s, "30-day launch checklist", "Owner Matrix")
left = [
    ("Aug 10", "Confirm launch date + sign-off on plan", ROSE),
    ("Aug 17", "Hire/assign marketing lead + content creator", ROSE),
    ("Aug 24", "Waitlist page live + RM30 offer configured", PINK),
    ("Aug 24", "GA4 + Meta events verified (10 events)", PINK),
    ("Aug 31", "Brevo sequences + UTM scheme live", PINK),
    ("Sep 21", "Media kit + press release drafted", PURPLE),
    ("Sep 21", "KOL wave 1 contracted (8–12)", PURPLE),
    ("Sep 21", "Launch pop-up venue + invite list", PURPLE),
    ("Sep 28", "30 anchor artists recruited (gate!)", PURPLE),
    ("Sep 28", "Paid campaigns live (test) → full Oct 5", PURPLE),
    ("Oct 5", "DAY 0 — official launch", DARK),
    ("Nov 4", "Referral live · wedding expo booked · Day-30 review", DARK),
    ("Dec 4", "Leish+ pilot to top clients", DARK),
    ("Jan 4", "Day-90 review + Q1'27 plan", DARK),
]
yy = 2.0
for date, task, c in left[:7]:
    box(s, 0.6, yy, 1.5, 0.52, fill=c, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.25)
    text(s, 0.6, yy + 0.09, 1.5, 0.36, date, size=11.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    text(s, 2.35, yy + 0.07, 3.6, 0.5, task, size=12, color=DARK, bold=True)
    yy += 0.6
yy = 2.0
for date, task, c in left[7:]:
    box(s, 6.6, yy, 1.5, 0.52, fill=c, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.25)
    text(s, 6.6, yy + 0.09, 1.5, 0.36, date, size=11.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    text(s, 8.35, yy + 0.07, 4.4, 0.5, task, size=12, color=DARK, bold=True)
    yy += 0.6

# ============================================================ 12. BLOCKERS
s = slide()
header(s, "Close by Aug 10", "Open Questions for the Founder")
qs = [
    ("1", "Confirm launch date", "Oct 5, 2026 proposed — first Monday after Hari Malaysia week, peak wedding season.", ROSE),
    ("2", "Budget sign-off", "RM 96,000 / 90 days, drawn from seed runway per financial forecast.", PINK),
    ("3", "Promo-code system", "Does admin support codes/credits, or is a dev task needed before Aug 24?", PURPLE),
    ("4", "Marketing owner", "In-house hire vs agency vs founder-led with this plan as the ops doc — #1 dependency.", RGBColor(0xD9, 0x77, 0x06)),
    ("5", "Brand assets", "IG/TikTok/FB handles (@leish.my?), logo kit, press photos.", RGBColor(0x0E, 0x94, 0x92)),
    ("6", "Leish+ timing", "Keep RM29/mo pilot internal for 90 days; public in Q1'27 (recommended).", RGBColor(0x7C, 0x3A, 0xED)),
    ("7", "Legal & compliance", "PDPA consent on waitlist capture + contest T&Cs before any giveaway.", RGBColor(0x65, 0xA3, 0x0D)),
]
yy = 2.0
for num, title, desc, c in qs:
    box(s, 0.6, yy, 0.62, 0.62, fill=c, shape=MSO_SHAPE.OVAL)
    text(s, 0.6, yy + 0.09, 0.62, 0.45, num, size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    text(s, 1.5, yy + 0.02, 4.3, 0.35, title, size=14, bold=True, color=DARK)
    text(s, 1.5, yy + 0.34, 11.2, 0.35, desc, size=11.5, color=GRAY)
    yy += 0.72

# ============================================================ 13. MEASUREMENT
s = slide()
header(s, "Weekly dashboard, every Monday", "Measurement & Attribution")
box(s, 0.6, 2.0, 5.9, 3.3, fill=LIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
text(s, 0.95, 2.3, 5.2, 0.35, "EVENTS TO VERIFY (GA4 + META)", size=13, color=ROSE, bold=True)
evs = ["waitlist_signup", "sign_up", "search", "view_artist", "booking_started", "booking_completed (value = booking value)", "payment_success", "artist_application_started", "share_link_clicked"]
yy = 2.75
for e in evs:
    box(s, 1.0, yy + 0.06, 0.08, 0.08, fill=ROSE, shape=MSO_SHAPE.OVAL)
    text(s, 1.2, yy, 5.2, 0.3, e, size=11.5, color=DARK, font="Consolas")
    yy += 0.285
box(s, 6.75, 2.0, 5.95, 3.3, fill=LIGHT2, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
text(s, 7.1, 2.3, 5.2, 0.35, "WEEKLY DASHBOARD SECTIONS", size=13, color=PURPLE, bold=True)
dash = [
    "Growth — waitlist, signups, artists, bookings, GMV, commission",
    "Funnel — search → artist view → booking_start → complete",
    "Paid — spend, ROAS, CAC, CPM/CPC by channel",
    "Quality — no-show, cancellation, avg rating, repeat rate",
    "Content — email list, opens/clicks, reach, top posts",
]
yy = 2.75
for d in dash:
    box(s, 7.15, yy + 0.06, 0.08, 0.08, fill=PURPLE, shape=MSO_SHAPE.OVAL)
    text(s, 7.35, yy, 5.3, 0.55, d, size=11.5, color=DARK, line_spacing=1.05)
    yy += 0.52
text(s, 0.6, 5.6, 12.1, 0.35, "ATTRIBUTION", size=13, color=GRAY, bold=True)
text(s, 0.6, 5.95, 12.1, 0.8, [
    [("UTM on every outbound link — source / medium / campaign / term / content (scheme in plan Appendix C). Capture source at signup. Report: bookings & CAC by channel.", {"size": 12, "color": GRAY})],
], line_spacing=1.15)
box(s, 0.6, 6.7, 12.1, 0.02, fill=RGBColor(0xE5, 0xE7, 0xEB))

# ============================================================ 14. CLOSING
s = slide(DARK2)
box(s, 0, 0, SW, 0.14, fill=ROSE)
box(s, 0, 7.36, SW, 0.14, fill=PURPLE)
box(s, 11.0, -1.7, 4.3, 4.3, fill=PINK, shape=MSO_SHAPE.OVAL)
box(s, -1.5, 5.2, 3.9, 3.9, fill=PURPLE, shape=MSO_SHAPE.OVAL)
text(s, 0.9, 2.2, 11.5, 1.2, [
    [("Ready to glow.", {"size": 46, "bold": True, "font": SERIF, "color": WHITE})],
], line_spacing=1.1)
text(s, 0.9, 3.5, 11.0, 0.5, [
    [("Full plan: ", {"size": 14, "color": MUTED}),
     ("docs/launch-marketing-plan.md", {"size": 14, "bold": True, "color": WHITE}),
     ("   ·   Summary: ", {"size": 14, "color": MUTED}),
     ("docs/launch-marketing-plan-summary.md", {"size": 14, "bold": True, "color": WHITE})],
])
text(s, 0.9, 4.15, 11.0, 0.5, [
    [("Next step: ", {"size": 14, "bold": True, "color": ROSE}),
     ("close the 7 blockers by Aug 10, then execute the 30-day checklist.", {"size": 14, "color": MUTED})],
])
text(s, 0.9, 6.3, 11.0, 0.5, [
    [("#BookBeautyAnywhere  ·  #NeverDMForADate  ·  Your Beauty, Perfected.", {"size": 13, "color": MUTED, "italic": True})],
])

prs.core_properties.title = "Leish! Launch Marketing Plan"
prs.core_properties.author = "Leish! (leish.my)"
prs.core_properties.subject = "90-day go-to-market plan — Oct 5, 2026"

os.makedirs(os.path.dirname(OUT), exist_ok=True)
prs.save(OUT)
print("Saved:", os.path.abspath(OUT), f"({os.path.getsize(OUT):,} bytes, {len(prs.slides._sldIdLst)} slides)")
