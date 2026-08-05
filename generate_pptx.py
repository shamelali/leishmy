#!/usr/bin/env python3
"""Generate Leish! Business Plan PowerPoint presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

OUTPUT_PATH = os.path.join(
    os.path.dirname(os.path.abspath(__file__)),
    "docs",
    "Leish_Business_Plan.pptx",
)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

DARK_PRIMARY = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_PINK = RGBColor(0xE8, 0x3E, 0x8C)
ACCENT_TEAL = RGBColor(0x20, 0xC9, 0x97)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF5)
DARK_TEXT = RGBColor(0x2D, 0x2D, 0x3F)
MEDIUM_TEXT = RGBColor(0x6C, 0x6C, 0x8A)
TABLE_HEADER_BG = RGBColor(0x1A, 0x1A, 0x2E)
TABLE_ALT_BG = RGBColor(0xF8, 0xF8, 0xFC)
BORDER_COLOR = RGBColor(0xDD, 0xDD, 0xE8)

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT


def add_bg(slide, color=DARK_PRIMARY):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_frame(slide, left, top, width, height, items, font_size=14, color=DARK_TEXT, line_spacing=1.3):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(4)
        p.level = 0
    return txBox


def add_accent_bar(slide, left, top, width, height, color=ACCENT_PINK):
    return add_shape_bg(slide, left, top, width, height, color)


def add_slide_number(slide, num, total):
    add_text_box(slide, Inches(12.2), Inches(7.05), Inches(1), Inches(0.35),
                 f"{num} / {total}", font_size=9, color=MEDIUM_TEXT, alignment=PP_ALIGN.RIGHT)


def add_section_header(slide, title, subtitle=None):
    add_accent_bar(slide, Inches(0.6), Inches(0.5), Inches(0.08), Inches(0.5), ACCENT_PINK)
    add_text_box(slide, Inches(0.85), Inches(0.45), Inches(11), Inches(0.6),
                 title, font_size=28, color=DARK_PRIMARY, bold=True)
    if subtitle:
        add_text_box(slide, Inches(0.85), Inches(1.0), Inches(11), Inches(0.4),
                     subtitle, font_size=14, color=MEDIUM_TEXT)


def add_table(slide, left, top, col_widths, headers, rows, font_size=11):
    total_w = sum(col_widths)
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, total_w, Inches(0.35 * (len(rows) + 1)))
    table = table_shape.table

    for ci, cw in enumerate(col_widths):
        table.columns[ci].width = cw

    for ci, header in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = header
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(font_size)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = "Calibri"
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER_BG

    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.color.rgb = DARK_TEXT
                p.font.name = "Calibri"
            if ri % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT_BG
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE

    return table_shape


TOTAL_SLIDES = 15

# ── Slide 1: Title ──
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK_PRIMARY)
add_shape_bg(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), ACCENT_PINK)
add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
             "Leish!", font_size=56, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(2.9), Inches(11), Inches(0.6),
             "Comprehensive Business Plan", font_size=24, color=RGBColor(0xCC, 0xCC, 0xDD), alignment=PP_ALIGN.CENTER)
add_accent_bar(slide, Inches(5.5), Inches(3.6), Inches(2.3), Inches(0.06), ACCENT_TEAL)
add_text_box(slide, Inches(1), Inches(4.0), Inches(11), Inches(0.5),
             "Malaysia's Leading Beauty & Wellness Marketplace", font_size=16, color=MEDIUM_TEXT, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.4),
             "August 2026  |  Confidential", font_size=12, color=MEDIUM_TEXT, alignment=PP_ALIGN.CENTER)
add_slide_number(slide, 1, TOTAL_SLIDES)

# ── Slide 2: Executive Summary ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, "Executive Summary", "Overview of Leish! and key targets")
add_bullet_frame(slide, Inches(0.85), Inches(1.6), Inches(5.5), Inches(5.0), [
    "• Malaysian beauty & wellness marketplace connecting customers with artists and studios",
    "• Built on Next.js 16 + Neon Postgres + Drizzle ORM",
    "• Payment via Billplz, email via Brevo, media via Cloudinary",
    "• Deployed on Vercel with Cloudflare edge infrastructure",
    "",
    "Key Targets (Month 12):",
    "  • 50,000 Monthly Active Users",
    "  • 2,000 Active Service Providers",
    "  • 15,000 Monthly Bookings",
    "  • RM 5M/month Gross Transaction Value",
], font_size=13)
add_table(slide, Inches(7.2), Inches(1.6),
          [Inches(2.2), Inches(1.5), Inches(1.5)],
          ["Metric", "Target", "Timeline"],
          [
              ["MAU", "50,000", "Month 12"],
              ["Providers", "2,000", "Month 12"],
              ["Monthly Bookings", "15,000", "Month 12"],
              ["GTV", "RM 5M/mo", "Month 12"],
              ["Revenue", "RM 135K/mo", "Month 12"],
              ["Break-Even", "Month 10-12", "—"],
          ], font_size=11)
add_slide_number(slide, 2, TOTAL_SLIDES)

# ── Slide 3: Company Overview ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, "Company Overview", "Mission, vision, and structure")
add_bullet_frame(slide, Inches(0.85), Inches(1.6), Inches(5.5), Inches(3.0), [
    "Mission: Democratize access to beauty and wellness services in Malaysia by making discovery, booking, and payment seamless for both consumers and providers.",
    "",
    "Vision: Become the leading beauty and wellness marketplace in Southeast Asia, starting with Malaysia.",
    "",
    "Legal: Malaysian Sdn Bhd (private limited company), registered in Kuala Lumpur.",
], font_size=13)
add_bullet_frame(slide, Inches(7.2), Inches(1.6), Inches(5.5), Inches(3.0), [
    "Founding Team:",
    "  • CEO / Founder — Strategy, fundraising, partnerships",
    "  • CTO — Architecture, engineering, infrastructure",
    "  • COO — Operations, onboarding, customer support",
    "  • CMO — Marketing, growth, brand",
    "",
    "Core Values:",
    "  • Customer-first design",
    "  • Local-first payments",
    "  • Transparent and secure",
    "  • Cloud-native and scalable",
], font_size=13)
add_slide_number(slide, 3, TOTAL_SLIDES)

# ── Slide 4: Market Analysis ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, "Market Analysis", "Industry size, target segments, and competition")
add_text_box(slide, Inches(0.85), Inches(1.5), Inches(11.5), Inches(0.4),
             "Industry: Malaysian beauty & wellness market ~RM 8B (2025), growing 12-15% CAGR", font_size=13, color=DARK_TEXT, bold=True)
add_table(slide, Inches(0.85), Inches(2.0),
          [Inches(2.5), Inches(4.5), Inches(2.5)],
          ["Segment", "Description", "Size Estimate"],
          [
              ["Urban Consumers (18-45)", "Tech-savvy Malaysians in KL, PJ, Penang, Johor", "~8M"],
              ["Beauty Artists", "Independent stylists, nail techs, aestheticians", "~50,000+"],
              ["Studios / Salons", "Small-medium salons seeking digital presence", "~10,000+"],
              ["Premium / Spa Clients", "High-income consumers for premium experiences", "~2M"],
          ], font_size=10)
add_text_box(slide, Inches(0.85), Inches(4.0), Inches(11.5), Inches(0.3),
             "Competitive Advantage", font_size=14, color=DARK_PRIMARY, bold=True)
add_bullet_frame(slide, Inches(0.85), Inches(4.35), Inches(11.5), Inches(2.5), [
    "1. Local-first payment integration (Billplz — FPX, credit card, e-wallet)",
    "2. Neon Auth — secure, modern authentication with session management",
    "3. Dashboard ecosystem — separate dashboards for admin, artists, and studios",
    "4. Cron-driven automation — booking reminders, payment reconciliation, lead follow-ups",
    "5. Cloud-native infrastructure — serverless, scalable, low operational overhead",
], font_size=11)
add_slide_number(slide, 4, TOTAL_SLIDES)

# ── Slide 5: Value Proposition ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, "Value Proposition", "What we offer to each stakeholder")
# Three columns
for ci, (title, items) in enumerate([
    ("For Customers", [
        "Discovery — browse by location, service, price, rating",
        "Convenience — book in minutes with real-time availability",
        "Security — Billplz escrow payment release",
        "Transparency — clear pricing, reviews, provider profiles",
        "Rewards — loyalty program with points and cashback",
    ]),
    ("For Providers", [
        "Client acquisition — steady stream of booked clients",
        "Business tools — calendar, staff mgmt, inventory, analytics",
        "Professional presence — portfolio, reviews, ratings",
        "Payment processing — automated collection & settlement",
        "Marketing — promotional tools, lead follow-up automation",
    ]),
    ("For the Platform", [
        "Commission revenue — % of each transaction",
        "Subscription revenue — premium plans for providers",
        "Advertising — sponsored listings and promoted profiles",
        "Lead generation — paid lead packages",
        "Data insights — aggregated market intelligence",
    ]),
]):
    left = Inches(0.5 + ci * 4.2)
    add_shape_bg(slide, left, Inches(1.5), Inches(3.9), Inches(0.45), ACCENT_PINK if ci == 0 else (ACCENT_TEAL if ci == 1 else RGBColor(0x6C, 0x63, 0xFF)))
    add_text_box(slide, left + Inches(0.15), Inches(1.52), Inches(3.6), Inches(0.4),
                 title, font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_bullet_frame(slide, left + Inches(0.1), Inches(2.1), Inches(3.7), Inches(4.5),
                     [f"• {item}" for item in items], font_size=11, color=DARK_TEXT)
add_slide_number(slide, 5, TOTAL_SLIDES)

# ── Slide 6: Business Model ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, "Business Model", "Revenue streams, pricing, and payment flow")
add_text_box(slide, Inches(0.85), Inches(1.5), Inches(5.5), Inches(0.3),
             "Revenue Streams", font_size=14, color=DARK_PRIMARY, bold=True)
add_table(slide, Inches(0.85), Inches(1.9),
          [Inches(2.8), Inches(2.2)],
          ["Stream", "Target % of Revenue"],
          [
              ["Transaction Commission (10-15%)", "60%"],
              ["Subscription Plans", "20%"],
              ["Advertising & Promotions", "10%"],
              ["Lead Generation", "5%"],
              ["Value-Added Services", "5%"],
          ], font_size=10)
add_text_box(slide, Inches(7.2), Inches(1.5), Inches(5.5), Inches(0.3),
             "Subscription Tiers", font_size=14, color=DARK_PRIMARY, bold=True)
add_table(slide, Inches(7.2), Inches(1.9),
          [Inches(1.3), Inches(1.5), Inches(2.7)],
          ["Tier", "Price (RM/mo)", "Features"],
          [
              ["Free", "RM 0", "Basic profile, 10 bookings/mo"],
              ["Pro", "RM 99", "Unlimited bookings, analytics, calendar"],
              ["Business", "RM 299", "Multi-staff, inventory, priority support"],
              ["Enterprise", "Custom", "White-label, API access, dedicated manager"],
          ], font_size=9)
add_text_box(slide, Inches(0.85), Inches(4.5), Inches(11.5), Inches(0.3),
             "Payment Flow", font_size=14, color=DARK_PRIMARY, bold=True)
add_bullet_frame(slide, Inches(0.85), Inches(4.85), Inches(11.5), Inches(2.0), [
    "1. Customer books a service and pays via Billplz",
    "2. Payment is held in escrow (Billplz pending)",
    "3. Service is completed",
    "4. Platform releases payment to provider (minus commission)",
    "5. Provider settles to bank account via FPX or transfer",
], font_size=11)
add_slide_number(slide, 6, TOTAL_SLIDES)

# ── Slide 7: Marketing Strategy ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, "Marketing Strategy", "Go-to-market phases and acquisition channels")
add_text_box(slide, Inches(0.85), Inches(1.5), Inches(11.5), Inches(0.3),
             "Go-to-Market Phases", font_size=14, color=DARK_PRIMARY, bold=True)
add_table(slide, Inches(0.85), Inches(1.9),
          [Inches(1.5), Inches(3.5), Inches(7.0)],
          ["Phase", "Timeline", "Key Activities"],
          [
              ["Seed", "Months 1-3", "Onboard 100 artists, 20 studios in KL; Instagram/TikTok campaigns; influencer collabs; referral program; Google Ads"],
              ["Scale", "Months 4-8", "Expand to Penang, Johor, Ipoh; launch subscriptions; content marketing; email campaigns; PR; beauty school partnerships"],
              ["Dominate", "Months 9-18", "All major cities; loyalty program; corporate booking; cross-border (Singapore, Indonesia); B2B product partnerships"],
          ], font_size=9)
add_text_box(slide, Inches(0.85), Inches(4.2), Inches(11.5), Inches(0.3),
             "Customer Acquisition Channels", font_size=14, color=DARK_PRIMARY, bold=True)
add_table(slide, Inches(0.85), Inches(4.55),
          [Inches(2.5), Inches(5.5), Inches(2.0)],
          ["Channel", "Strategy", "Target CPA (RM)"],
          [
              ["Instagram/TikTok", "Reels, influencer collabs, paid ads", "RM 15-25"],
              ["Google Ads", "Search + Display, retargeting", "RM 20-30"],
              ["Referral Program", "RM 20 credit per successful referral", "RM 10-15"],
              ["Content Marketing", "SEO blog, YouTube tutorials", "RM 5-10"],
              ["Email Marketing", "Brevo campaigns, welcome sequences", "RM 3-8"],
              ["WhatsApp Business", "Webhook notifications, broadcasts", "RM 2-5"],
          ], font_size=9)
add_slide_number(slide, 7, TOTAL_SLIDES)

# ── Slide 8: Technology Stack ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, "Technology Stack", "Architecture and infrastructure")
add_table(slide, Inches(0.85), Inches(1.5),
          [Inches(2.5), Inches(9.0)],
          ["Layer", "Technology"],
          [
              ["Frontend", "Next.js 16 (App Router), TypeScript, Tailwind CSS 4"],
              ["Backend", "Next.js API Routes, Drizzle ORM"],
              ["Database", "Neon (serverless Postgres)"],
              ["Auth", "@neondatabase/auth (Neon Auth / Better Auth)"],
              ["Payments", "Billplz (Malaysian payment gateway)"],
              ["Email", "Brevo (@getbrevo/brevo)"],
              ["Storage", "Cloudinary (images/media)"],
              ["CDN/Edge", "Cloudflare (Workers, KV, R2)"],
              ["Monitoring", "Sentry (error tracking, performance)"],
              ["Hosting", "Vercel"],
              ["Testing", "Playwright (E2E)"],
              ["CI/CD", "GitHub Actions (auto-deploy on push to main)"],
          ], font_size=10)
add_text_box(slide, Inches(0.85), Inches(5.9), Inches(11.5), Inches(0.3),
             "Key Infrastructure Components", font_size=14, color=DARK_PRIMARY, bold=True)
add_bullet_frame(slide, Inches(0.85), Inches(6.2), Inches(5.5), Inches(1.0), [
    "Middleware (src/proxy.ts): Dashboard auth, rate limiting, CSP headers",
], font_size=10)
add_bullet_frame(slide, Inches(7.0), Inches(6.2), Inches(5.5), Inches(1.0), [
    "Cron Jobs: Sync auth, reconcile payments, booking reminders, lead follow-ups",
], font_size=10)
add_slide_number(slide, 8, TOTAL_SLIDES)

# ── Slide 9: Financial Projections ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, "Financial Projections", "Startup costs, operating costs, and revenue forecast")
add_text_box(slide, Inches(0.85), Inches(1.5), Inches(5.5), Inches(0.3),
             "Startup Costs (One-Time)", font_size=13, color=DARK_PRIMARY, bold=True)
add_table(slide, Inches(0.85), Inches(1.85),
          [Inches(3.5), Inches(1.5)],
          ["Item", "Cost (RM)"],
          [
              ["Company registration & legal", "5,000"],
              ["Brand identity & design", "10,000"],
              ["Initial development (6 months)", "150,000"],
              ["Marketing launch budget", "30,000"],
              ["Infrastructure setup", "10,000"],
              ["Miscellaneous", "5,000"],
              ["TOTAL", "210,000"],
          ], font_size=9)
add_text_box(slide, Inches(7.2), Inches(1.5), Inches(5.5), Inches(0.3),
             "Monthly Operating Costs", font_size=13, color=DARK_PRIMARY, bold=True)
add_table(slide, Inches(7.2), Inches(1.85),
          [Inches(3.5), Inches(1.5)],
          ["Item", "Cost (RM)"],
          [
              ["Vercel hosting", "500"],
              ["Neon database", "1,000"],
              ["Cloudinary storage/bandwidth", "800"],
              ["Billplz transaction fees", "Variable"],
              ["Brevo email", "200"],
              ["Cloudflare Workers/KV", "300"],
              ["Upstash Redis", "200"],
              ["Sentry", "150"],
              ["Marketing spend", "15,000"],
              ["Staff salaries (4 FTE)", "40,000"],
              ["Office / co-working", "2,000"],
              ["Miscellaneous", "1,000"],
              ["TOTAL (fixed)", "~61,150 + variable"],
          ], font_size=9)
add_text_box(slide, Inches(0.85), Inches(4.8), Inches(11.5), Inches(0.3),
             "Revenue Projections", font_size=13, color=DARK_PRIMARY, bold=True)
add_table(slide, Inches(0.85), Inches(5.15),
          [Inches(1.2), Inches(1.8), Inches(1.8), Inches(2.0), Inches(1.8), Inches(1.5), Inches(1.8)],
          ["Month", "Bookings", "GTV", "Commission (10%)", "Subscriptions", "Other", "Total Rev"],
          [
              ["M1", "500", "25,000", "2,500", "0", "0", "2,500"],
              ["M3", "2,000", "100,000", "10,000", "2,000", "500", "12,500"],
              ["M6", "8,000", "400,000", "40,000", "10,000", "2,000", "52,000"],
              ["M9", "12,000", "600,000", "60,000", "25,000", "5,000", "90,000"],
              ["M12", "15,000", "750,000", "75,000", "50,000", "10,000", "135,000"],
          ], font_size=9)
add_slide_number(slide, 9, TOTAL_SLIDES)

# ── Slide 10: Break-Even & Funding ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, "Break-Even Analysis & Funding", "Path to profitability and capital requirements")
add_text_box(slide, Inches(0.85), Inches(1.6), Inches(5.5), Inches(0.3),
             "Break-Even Analysis", font_size=14, color=DARK_PRIMARY, bold=True)
add_bullet_frame(slide, Inches(0.85), Inches(2.0), Inches(5.5), Inches(2.5), [
    "• Fixed monthly costs: ~RM 61,150",
    "• Break-even GTV (at 10% commission): RM 611,500/month",
    "• Projected to reach break-even by Month 10-12",
    "• Revenue grows from RM 2,500 (M1) to RM 135,000 (M12)",
    "• Commission rate: 10% standard, 8% for subscribed providers",
], font_size=12)
add_text_box(slide, Inches(7.2), Inches(1.6), Inches(5.5), Inches(0.3),
             "Funding Requirements", font_size=14, color=DARK_PRIMARY, bold=True)
add_table(slide, Inches(7.2), Inches(2.0),
          [Inches(2.5), Inches(1.5)],
          ["Use of Funds", "Allocation"],
          [
              ["Product dev & engineering", "40%"],
              ["Marketing & user acquisition", "30%"],
              ["Operations & infrastructure", "15%"],
              ["Legal, admin & compliance", "10%"],
              ["Contingency", "5%"],
          ], font_size=10)
add_text_box(slide, Inches(7.2), Inches(4.2), Inches(5.5), Inches(0.3),
             "Seed Round: RM 500,000", font_size=14, color=ACCENT_PINK, bold=True)
add_bullet_frame(slide, Inches(7.2), Inches(4.6), Inches(5.5), Inches(1.5), [
    "• Covers 18 months of operations + development",
    "• Target: close by Month 6",
    "• Leads to Series A readiness by Month 18",
], font_size=11)
add_text_box(slide, Inches(0.85), Inches(5.2), Inches(11.5), Inches(0.3),
             "Key Financial Metrics", font_size=14, color=DARK_PRIMARY, bold=True)
add_table(slide, Inches(0.85), Inches(5.55),
          [Inches(3.0), Inches(2.5), Inches(2.5), Inches(3.0)],
          ["Metric", "Target", "Rationale", "Measurement"],
          [
              ["CAC", "< RM 25", "Affordable digital acquisition", "Total marketing spend / new customers"],
              ["LTV", "> RM 200", "Strong retention & repeat bookings", "Avg. bookings per customer × avg. value"],
              ["LTV:CAC", "> 8:1", "Healthy unit economics", "LTV / CAC"],
              ["Monthly Churn", "< 5%", "Strong provider & customer retention", "Lost customers / total customers"],
          ], font_size=9)
add_slide_number(slide, 10, TOTAL_SLIDES)

# ── Slide 11: Growth Strategy ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, "Growth Strategy & Product Roadmap", "From MVP to regional leader")
roadmap = [
    ("Q3 2026 (Now)", "MVP Launch", [
        "Core booking flow live",
        "Onboard 100+ artists in KL",
        "Basic search and filtering",
        "Billplz payment integration",
    ]),
    ("Q4 2026", "Product Enhancement", [
        "Mobile PWA optimization",
        "Review and rating system",
        "Loyalty points program",
        "Advanced search filters",
    ]),
    ("Q1-Q2 2027", "Scale & Expand", [
        "Mobile native app (React Native/Flutter)",
        "AI-powered recommendations",
        "Multi-city expansion (Penang, Johor, Ipoh)",
        "Subscription plans for providers",
    ]),
    ("Q3-Q4 2027", "Market Leadership", [
        "Cross-border expansion (Singapore)",
        "B2B marketplace (product suppliers)",
        "Insurance integration for providers",
        "API marketplace & white-label offering",
    ]),
]
y = Inches(1.5)
for quarter, phase, items in roadmap:
    add_shape_bg(slide, Inches(0.85), y, Inches(0.08), Inches(0.35), ACCENT_PINK)
    add_text_box(slide, Inches(1.1), y, Inches(2.5), Inches(0.35),
                 quarter, font_size=11, color=ACCENT_PINK, bold=True)
    add_text_box(slide, Inches(3.5), y, Inches(2.0), Inches(0.35),
                 phase, font_size=11, color=DARK_PRIMARY, bold=True)
    add_bullet_frame(slide, Inches(5.5), y, Inches(7.0), Inches(0.35),
                     [f"• {item}" for item in items], font_size=10, color=DARK_TEXT)
    y += Inches(0.55)
add_text_box(slide, Inches(0.85), Inches(4.2), Inches(11.5), Inches(0.3),
             "Key Growth Metrics", font_size=14, color=DARK_PRIMARY, bold=True)
add_table(slide, Inches(0.85), Inches(4.55),
          [Inches(3.5), Inches(2.5), Inches(2.5), Inches(3.0)],
          ["Metric", "Target", "Why It Matters", "How We Measure"],
          [
              ["CAC", "< RM 25", "Affordable acquisition", "Marketing spend / new customers"],
              ["LTV", "> RM 200", "Strong customer value", "Avg bookings × avg value"],
              ["LTV:CAC", "> 8:1", "Healthy unit economics", "LTV / CAC ratio"],
              ["Churn Rate", "< 5%", "Strong retention", "Lost customers / total"],
              ["NPS", "> 50", "Customer satisfaction", "Survey-based scoring"],
          ], font_size=9)
add_slide_number(slide, 11, TOTAL_SLIDES)

# ── Slide 12: Risk Analysis ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, "Risk Analysis", "Key risks and mitigation strategies")
add_table(slide, Inches(0.5), Inches(1.5),
          [Inches(2.8), Inches(1.2), Inches(1.2), Inches(7.0)],
          ["Risk", "Likelihood", "Impact", "Mitigation"],
          [
              ["Low initial provider supply", "High", "High", "Aggressive onboarding incentives, beauty school partnerships"],
              ["Payment fraud / disputes", "Medium", "High", "Billplz escrow, KYC verification, fraud detection rules"],
              ["Competitor entry", "Medium", "Medium", "First-mover advantage, local payment integration, network effects"],
              ["Regulatory changes", "Low", "Medium", "Legal counsel, compliance monitoring, flexible structure"],
              ["Technology failure", "Low", "High", "Multi-region Neon, Vercel edge, automated backups, Sentry monitoring"],
              ["Provider churn", "Medium", "Medium", "Subscription lock-in, analytics tools, marketing support"],
              ["Customer trust", "Medium", "High", "Verified profiles, secure payments, transparent reviews"],
              ["Cash flow gap", "Medium", "High", "Lean operations, milestone-based funding, revenue diversification"],
          ], font_size=9)
add_slide_number(slide, 12, TOTAL_SLIDES)

# ── Slide 13: Team & Milestones ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, "Team, Hiring Plan & Milestones", "People and key milestones")
add_text_box(slide, Inches(0.85), Inches(1.5), Inches(5.5), Inches(0.3),
             "Current Team (4 FTE)", font_size=13, color=DARK_PRIMARY, bold=True)
add_table(slide, Inches(0.85), Inches(1.85),
          [Inches(1.5), Inches(4.0)],
          ["Role", "Responsibilities"],
              [
              ["CEO / Founder", "Strategy, fundraising, partnerships"],
              ["CTO", "Architecture, engineering, infrastructure"],
              ["COO", "Operations, onboarding, customer support"],
              ["CMO", "Marketing, growth, brand"],
          ], font_size=9)
add_text_box(slide, Inches(7.2), Inches(1.5), Inches(5.5), Inches(0.3),
             "Hiring Plan", font_size=13, color=DARK_PRIMARY, bold=True)
add_table(slide, Inches(7.2), Inches(1.85),
          [Inches(1.2), Inches(2.5), Inches(2.0)],
          ["Month", "Role", "Purpose"],
          [
              ["M3", "2x Full-Stack Engineers", "Scale product development"],
              ["M4", "1x Marketing Specialist", "Execute growth campaigns"],
              ["M6", "1x Customer Support Agent", "Handle growing user base"],
              ["M9", "1x Data Analyst", "Growth analytics & optimization"],
              ["M12", "1x Business Development", "Partnerships & expansion"],
          ], font_size=9)
add_text_box(slide, Inches(0.85), Inches(4.0), Inches(11.5), Inches(0.3),
             "Key Milestones", font_size=13, color=DARK_PRIMARY, bold=True)
add_table(slide, Inches(0.85), Inches(4.35),
          [Inches(3.0), Inches(1.5), Inches(7.0)],
          ["Milestone", "Target", "Success Criteria"],
          [
              ["MVP Launch", "Month 1", "Core booking flow live, 50+ artists onboarded"],
              ["First 1,000 Users", "Month 2", "1,000 registered customers"],
              ["First 100 Bookings", "Month 2", "100 completed bookings"],
              ["Payment Integration Live", "Month 1", "Billplz fully integrated and tested"],
              ["500 Active Providers", "Month 6", "500+ artists and studios on platform"],
              ["RM 100K Monthly Revenue", "Month 9", "Sustained monthly revenue > RM 100K"],
              ["Break-Even", "Month 10-12", "Monthly revenue exceeds monthly costs"],
              ["Multi-City Expansion", "Month 12", "Operations in 3+ cities"],
              ["Seed Funding Closed", "Month 6", "RM 500K raised"],
              ["Series A Readiness", "Month 18", "Strong metrics, ready for Series A"],
          ], font_size=8)
add_slide_number(slide, 13, TOTAL_SLIDES)

# ── Slide 14: KPIs & Exit Strategy ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, "KPIs & Exit Strategy", "How we measure success and future paths")
add_text_box(slide, Inches(0.85), Inches(1.5), Inches(5.5), Inches(0.3),
             "Key Performance Indicators", font_size=13, color=DARK_PRIMARY, bold=True)
add_table(slide, Inches(0.85), Inches(1.85),
          [Inches(2.5), Inches(1.5), Inches(1.5)],
          ["Category", "KPI", "Target"],
          [
              ["Customer", "MAU", "50,000 (M12)"],
              ["Customer", "CAC", "< RM 25"],
              ["Customer", "LTV", "> RM 200"],
              ["Customer", "Booking Conversion", "> 15%"],
              ["Customer", "Repeat Booking Rate", "> 40%"],
              ["Provider", "Active Providers", "2,000 (M12)"],
              ["Provider", "Provider Churn", "< 5%/month"],
              ["Financial", "GTV", "RM 5M/month (M12)"],
              ["Financial", "Gross Margin", "> 60%"],
              ["Operational", "System Uptime", "99.9%"],
              ["Operational", "Page Load Time", "< 2 seconds"],
          ], font_size=9)
add_text_box(slide, Inches(7.2), Inches(1.5), Inches(5.5), Inches(0.3),
             "Exit Strategy", font_size=13, color=DARK_PRIMARY, bold=True)
add_bullet_frame(slide, Inches(7.2), Inches(1.9), Inches(5.5), Inches(2.5), [
    "1. Acquisition — By a larger beauty/wellness platform (e.g., Style Theory, Zalora, Fresha, Booksy)",
    "2. IPO — Long-term goal of listing on Malaysian stock exchange (MYX)",
    "3. Strategic Partnership — Acquisition of technology/IP by a major player in beauty/retail",
    "",
    "Target exit timeline: 5-7 years from launch",
    "Series A readiness target: Month 18",
], font_size=11)
add_text_box(slide, Inches(0.85), Inches(5.5), Inches(11.5), Inches(0.3),
             "Operational Targets", font_size=13, color=DARK_PRIMARY, bold=True)
add_table(slide, Inches(0.85), Inches(5.85),
          [Inches(3.0), Inches(2.0), Inches(2.0), Inches(3.5)],
          ["Metric", "Target", "Current", "Notes"],
          [
              ["API Response Time", "< 200ms", "TBD", "Neon edge + Vercel"],
              ["Error Rate", "< 0.1%", "TBD", "Sentry monitoring"],
              ["Support Response", "< 4 hours", "TBD", "Business hours"],
              ["Email Delivery", "> 99%", "TBD", "Brevo + Cloudflare routing"],
          ], font_size=9)
add_slide_number(slide, 14, TOTAL_SLIDES)

# ── Slide 15: Thank You ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_PRIMARY)
add_shape_bg(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), ACCENT_PINK)
add_text_box(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.0),
             "Thank You", font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_accent_bar(slide, Inches(5.5), Inches(3.1), Inches(2.3), Inches(0.06), ACCENT_TEAL)
add_text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(0.5),
             "Leish! — Malaysia's Beauty & Wellness Marketplace", font_size=18, color=RGBColor(0xCC, 0xCC, 0xDD), alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.4),
             "leish.my  |  Confidential  |  August 2026", font_size=14, color=MEDIUM_TEXT, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.4),
             "Questions? Contact the Leish! founding team", font_size=14, color=MEDIUM_TEXT, alignment=PP_ALIGN.CENTER)
add_slide_number(slide, 15, TOTAL_SLIDES)

# ── Save ──
os.makedirs(os.path.dirname(OUTPUT_PATH), exist_ok=True)
prs.save(OUTPUT_PATH)
print(f"Presentation saved to: {OUTPUT_PATH}")
print(f"Slides: {len(prs.slides)}")